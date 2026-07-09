#!/usr/bin/env python3
"""Generate synthetic Office document fixtures and .eml files for HIVE tests."""

from __future__ import annotations

import base64
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

SAMPLES_DIR = Path(__file__).parent / "samples"

DOCX_PATH = SAMPLES_DIR / "test_attachment.docx"
XLSX_PATH = SAMPLES_DIR / "test_attachment.xlsx"
PPTX_PATH = SAMPLES_DIR / "test_attachment.pptx"

DOCX_EML_PATH = SAMPLES_DIR / "docx_attachment.eml"
XLSX_EML_PATH = SAMPLES_DIR / "xlsx_attachment.eml"
PPTX_EML_PATH = SAMPLES_DIR / "pptx_attachment.eml"


def create_docx(path: Path) -> bytes:
    """Create a synthetic DOCX containing paragraphs and a table with URLs."""
    document = Document()
    document.add_paragraph("Please review the following document.")
    document.add_paragraph(
        "Approval link: https://malicious-docx-link.com/approve?id=1234"
    )
    document.add_paragraph("For support visit: http://docx-support.net/help")

    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Item"
    table.rows[0].cells[1].text = "Link"
    table.rows[1].cells[0].text = "Portal"
    table.rows[1].cells[1].text = "https://docx-table-link.com/portal"

    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)
    return path.read_bytes()


def create_xlsx(path: Path) -> bytes:
    """Create a synthetic XLSX workbook with URL cells."""
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet["A1"] = "Description"
    sheet["B1"] = "URL"
    sheet["A2"] = "Malicious link"
    sheet["B2"] = "https://malicious-xlsx-link.com/payload?ref=9921"
    sheet["A3"] = "Support"
    sheet["B3"] = "http://xlsx-support.net/contact"

    path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(path)
    return path.read_bytes()


def create_pptx(path: Path) -> bytes:
    """Create a synthetic PPTX presentation with URLs in a text box."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    text_frame = textbox.text_frame
    text_frame.text = (
        "Click here to verify: https://malicious-pptx-link.com/verify?token=abc"
    )
    support_paragraph = text_frame.add_paragraph()
    support_paragraph.text = "Support: http://pptx-support.net/help"

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return path.read_bytes()


def _wrap_base64(data: bytes) -> str:
    """Encode bytes as base64 with RFC 2045 line wrapping."""
    encoded = base64.b64encode(data).decode("ascii")
    return "\r\n".join(encoded[index : index + 76] for index in range(0, len(encoded), 76))


def create_eml(
    *,
    attachment_bytes: bytes,
    content_type: str,
    attachment_filename: str,
    subject: str,
    body_url: str,
    boundary: str,
    message_id: str,
    body_intro: str,
) -> str:
    """Build a synthetic phishing .eml with a base64 Office attachment."""
    wrapped_attachment = _wrap_base64(attachment_bytes)

    return f"""Return-Path: <reports@quarterly-docs.com>
Received: from mail.quarterly-docs.com (mail.quarterly-docs.com [198.51.100.44])
        by mx1.protection.outlook.com with SMTP id o1si9932456qkd
        for <finance@vunhst.nhs.uk>; Wed, 13 Mar 2024 10:00:04 +0000
Received: from localhost (localhost [127.0.0.1])
        by mail.quarterly-docs.com with ESMTP id s2t3u4v5w6
        for <finance@vunhst.nhs.uk>; Wed, 13 Mar 2024 10:00:00 +0000
Received-SPF: fail (domain of reports@quarterly-docs.com does not designate
        198.51.100.44 as permitted sender)
        receiver=mx1.protection.outlook.com;
        client-ip=198.51.100.44;
        helo=mail.quarterly-docs.com
Authentication-Results: mx1.protection.outlook.com;
        dmarc=fail action=none header.from=quarterly-docs.com;
        spf=fail smtp.mailfrom=reports@quarterly-docs.com;
        dkim=none
From: "Quarterly Reports" <reports@quarterly-docs.com>
Reply-To: reports@quarterly-docs.com
To: finance@vunhst.nhs.uk
Subject: {subject}
Date: Wed, 13 Mar 2024 10:00:00 +0000
Message-ID: <{message_id}>
MIME-Version: 1.0
X-Mailer: SendGrid
X-Originating-IP: 198.51.100.44
Content-Type: multipart/mixed; boundary="{boundary}"

--{boundary}
Content-Type: text/plain; charset="UTF-8"
Content-Transfer-Encoding: 7bit

Dear Finance Team,

{body_intro}

If the attachment does not open, use this link:
{body_url}

Regards,
Quarterly Reports Team

--{boundary}
Content-Type: {content_type}
Content-Disposition: attachment; filename="{attachment_filename}"
Content-Transfer-Encoding: base64

{wrapped_attachment}

--{boundary}--
"""


def _write_file(path: Path, content: str | bytes) -> None:
    """Write text or bytes to disk, creating parent directories as needed."""
    path.parent.mkdir(parents=True, exist_ok=True)
    if isinstance(content, bytes):
        path.write_bytes(content)
    else:
        path.write_text(content, encoding="utf-8", newline="\n")


def main() -> None:
    """Generate Office document fixtures and matching .eml files."""
    SAMPLES_DIR.mkdir(parents=True, exist_ok=True)

    docx_bytes = create_docx(DOCX_PATH)
    _write_file(
        DOCX_EML_PATH,
        create_eml(
            attachment_bytes=docx_bytes,
            content_type=(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
            attachment_filename="report.docx",
            subject="Q1 Word Report - Review Required",
            body_url="https://quarterly-docs.com/docx-report?ref=Q12024",
            boundary="=_HIVE_DOCX_BOUNDARY_001",
            message_id="20240313100000.docx8821345@mail.quarterly-docs.com",
            body_intro="Please find attached the Q1 Word report for your review.",
        ),
    )

    xlsx_bytes = create_xlsx(XLSX_PATH)
    _write_file(
        XLSX_EML_PATH,
        create_eml(
            attachment_bytes=xlsx_bytes,
            content_type=(
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ),
            attachment_filename="data.xlsx",
            subject="Q1 Excel Report - Review Required",
            body_url="https://quarterly-docs.com/xlsx-report?ref=Q12024",
            boundary="=_HIVE_XLSX_BOUNDARY_001",
            message_id="20240313100001.xlsx8821346@mail.quarterly-docs.com",
            body_intro="Please find attached the Q1 Excel report for your review.",
        ),
    )

    pptx_bytes = create_pptx(PPTX_PATH)
    _write_file(
        PPTX_EML_PATH,
        create_eml(
            attachment_bytes=pptx_bytes,
            content_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            attachment_filename="slides.pptx",
            subject="Q1 PowerPoint Report - Review Required",
            body_url="https://quarterly-docs.com/pptx-report?ref=Q12024",
            boundary="=_HIVE_PPTX_BOUNDARY_001",
            message_id="20240313100002.pptx8821347@mail.quarterly-docs.com",
            body_intro="Please find attached the Q1 PowerPoint report for your review.",
        ),
    )

    created_files = [
        DOCX_PATH,
        XLSX_PATH,
        PPTX_PATH,
        DOCX_EML_PATH,
        XLSX_EML_PATH,
        PPTX_EML_PATH,
    ]
    for path in created_files:
        print(f"Created {path} ({path.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
