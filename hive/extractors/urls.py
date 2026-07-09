"""URL extraction, deduplication, and defanging for HIVE forensic output.

Extracts scheme-prefixed URLs from email bodies, HTML, attachments, and
nested messages. All findings are defanged before output for safe analyst
review in urls.txt.
"""

from __future__ import annotations

import logging
import re
import unicodedata
from dataclasses import dataclass
from io import BytesIO
from urllib.parse import urlparse

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from striprtf.striprtf import rtf_to_text

from hive.extractors.headers import defang
from hive.parser.common import Attachment, ParsedEmail

logger = logging.getLogger(__name__)

# Scheme-prefixed URLs in plain text (http/https/ftp)
_URL_RE = re.compile(
    r"https?://[^\s<>\"'\]\[}{)(\|\\]+|ftp://[^\s<>\"'\]\[}{)(\|\\]+",
    re.IGNORECASE,
)

# href and src attribute values with http/https/ftp schemes
_HREF_SRC_RE = re.compile(
    r'(?:href|src)\s*=\s*(["\'])(https?://[^"\']+|ftp://[^"\']+)\1',
    re.IGNORECASE,
)

# Trailing punctuation often captured at end of sentences
_TRAILING_PUNCTUATION = ".,;:!?)'\""

_URL_SHORTENERS = frozenset({
    "bit.ly", "tinyurl.com", "t.co", "goo.gl", "ow.ly",
    "rebrand.ly", "short.io", "tiny.cc", "is.gd", "buff.ly",
    "cutt.ly", "bl.ink", "rb.gy", "shorte.st", "adf.ly",
    "bc.vc", "trib.al", "snip.ly", "clck.ru", "qr.ae",
    "x.co", "v.gd", "qurl.com", "budurl.com", "fur.ly",
    "lnkd.in", "youtu.be", "amzn.to", "fb.me", "po.st",
    "dlvr.it", "ift.tt", "soo.gd", "s2r.co", "clk.im",
    "tny.im", "tr.im", "url4.eu", "ff.im", "su.pr",
})

_SCRIPT_MARKERS = ("LATIN", "CYRILLIC", "GREEK", "ARABIC", "ARMENIAN", "GEORGIAN")


@dataclass
class UrlFinding:
    """A single defanged URL with its source location."""

    raw_url: str  # original URL before defanging
    defanged_url: str  # defanged version safe for analyst output
    source: str  # human-readable source label, e.g. body:plain or attachment:file.pdf
    depth: int  # email nesting depth where the URL was found (0 = top level)
    is_shortener: bool = False
    is_punycode: bool = False
    homoglyph_detail: str = ""


def _is_url_shortener(url: str) -> bool:
    """Return True if the URL uses a known shortener domain."""
    try:
        host = urlparse(url).hostname or ""
        host = host.lower().lstrip("www.")
        return host in _URL_SHORTENERS
    except Exception:
        return False


def check_punycode(url: str) -> tuple[bool, str]:
    """
    Detect punycode-encoded domains in a URL.

    Returns (is_punycode, decoded_domain). decoded_domain is the human-readable
    form if punycode was found, empty string otherwise. Never raises.
    """
    try:
        host = urlparse(url).hostname or ""
        if not host:
            return False, ""
        if "xn--" not in host.lower():
            return False, ""
        try:
            decoded = host.encode("ascii").decode("idna")
            if decoded != host:
                return True, decoded
        except (UnicodeError, UnicodeDecodeError):
            return True, host
        return False, ""
    except Exception:
        return False, ""


def _script_from_char(char: str) -> str | None:
    """Return a script marker for a hostname character, or None for separators."""
    if char in ".-":
        return None
    if char.isascii():
        return "LATIN"
    try:
        name = unicodedata.name(char, "")
        for marker in _SCRIPT_MARKERS:
            if marker in name:
                return marker
    except Exception:
        pass
    return "OTHER"


def _check_homoglyphs(url: str) -> tuple[bool, str]:
    """
    Detect non-ASCII characters and mixed Unicode scripts in a domain.

    Returns (is_suspicious, detail_string). Never raises.
    """
    try:
        host = urlparse(url).hostname or ""
        if not host or host.isascii():
            return False, ""

        scripts: set[str] = set()
        suspicious_chars: list[str] = []

        for char in host:
            script = _script_from_char(char)
            if script is None:
                continue
            scripts.add(script)
            if not char.isascii():
                try:
                    char_name = unicodedata.name(char, "UNKNOWN")
                    code_point = f"U+{ord(char):04X}"
                    suspicious_chars.append(f"{char} ({code_point}, {char_name})")
                except Exception:
                    suspicious_chars.append(char)

        if len(scripts) > 1:
            detail_base = "mixed-script domain (possible homoglyph attack)"
        else:
            detail_base = "non-ASCII characters in domain"

        if suspicious_chars:
            detail = f"{detail_base}\n{', '.join(suspicious_chars)}"
        else:
            detail = detail_base

        return True, detail
    except Exception:
        return False, ""


def get_url_warnings(findings: list[UrlFinding]) -> list[str]:
    """
    Return warning strings for flagged URLs for summary output.

    One warning per flagged URL per category. Never raises.
    """
    warnings: list[str] = []
    try:
        for finding in findings:
            if finding.is_shortener:
                warnings.append(
                    "URL shortener detected — destination unknown: "
                    f"{finding.defanged_url}"
                )
            if finding.is_punycode:
                _, decoded = check_punycode(finding.raw_url)
                encoded_host = urlparse(finding.raw_url).hostname or ""
                encoded_defanged = defang(encoded_host) if encoded_host else finding.defanged_url
                decoded_defanged = defang(decoded) if decoded else ""
                warnings.append(
                    "Punycode domain detected — renders as: "
                    f"{encoded_defanged} → {decoded_defanged}"
                )
            if finding.homoglyph_detail:
                main_detail = finding.homoglyph_detail.split("\n", 1)[0]
                warnings.append(
                    f"Suspicious Unicode in domain — {main_detail}: "
                    f"{finding.defanged_url}"
                )
    except Exception:
        logger.exception("Failed to build URL warnings")
    return warnings


def _strip_trailing_punctuation(url: str) -> str:
    """Remove trailing sentence punctuation from a URL match."""
    while url and url[-1] in _TRAILING_PUNCTUATION:
        url = url[:-1]
    return url


def _dedupe_preserve_order(urls: list[str]) -> list[str]:
    """Deduplicate URLs within a single source while preserving order."""
    seen: set[str] = set()
    unique: list[str] = []
    for url in urls:
        key = url.lower()
        if key in seen:
            continue
        seen.add(key)
        unique.append(url)
    return unique


def _extract_raw_urls(text: str) -> list[str]:
    """Apply the raw URL regex and strip trailing punctuation."""
    try:
        if not text:
            return []
        urls: list[str] = []
        for match in _URL_RE.finditer(text):
            url = _strip_trailing_punctuation(match.group())
            if url:
                urls.append(url)
        return urls
    except Exception:
        logger.exception("Failed to extract raw URLs from text")
        return []


def _extract_href_urls(html: str) -> list[str]:
    """Extract href and src attribute values that use http/https/ftp."""
    try:
        if not html:
            return []
        urls: list[str] = []
        for match in _HREF_SRC_RE.finditer(html):
            url = _strip_trailing_punctuation(match.group(2))
            if url:
                urls.append(url)
        return urls
    except Exception:
        logger.exception("Failed to extract href/src URLs from HTML")
        return []


def _extract_from_html_content(html: str) -> list[str]:
    """Extract URLs from HTML using href/src attributes and inline text."""
    combined = _extract_href_urls(html) + _extract_raw_urls(html)
    return _dedupe_preserve_order(combined)


def _extract_from_pdf(data: bytes) -> list[str]:
    """Extract URLs from PDF attachment text using pypdf."""
    try:
        reader = PdfReader(BytesIO(data))
        page_texts: list[str] = []
        for page in reader.pages:
            page_texts.append(page.extract_text() or "")
        return _extract_raw_urls("\n".join(page_texts))
    except Exception:
        logger.warning("Failed to extract URLs from PDF attachment", exc_info=True)
        return []


def _extract_from_docx(data: bytes) -> list[str]:
    """Extract URLs from DOCX paragraphs and table cells."""
    try:
        document = Document(BytesIO(data))
        texts: list[str] = []
        for paragraph in document.paragraphs:
            texts.append(paragraph.text)
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
        return _extract_raw_urls("\n".join(texts))
    except Exception:
        logger.warning("Failed to extract URLs from DOCX attachment", exc_info=True)
        return []


def _extract_from_xlsx(data: bytes) -> list[str]:
    """Extract URLs from XLSX workbook sheets, rows, and cells."""
    try:
        workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
        texts: list[str] = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        texts.append(str(cell.value))
        workbook.close()
        return _extract_raw_urls("\n".join(texts))
    except Exception:
        logger.warning("Failed to extract URLs from XLSX attachment", exc_info=True)
        return []


def _extract_from_pptx(data: bytes) -> list[str]:
    """Extract URLs from PPTX slide shape text frames."""
    try:
        presentation = Presentation(BytesIO(data))
        texts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    texts.append(paragraph.text)
        return _extract_raw_urls("\n".join(texts))
    except Exception:
        logger.warning("Failed to extract URLs from PPTX attachment", exc_info=True)
        return []


def _extract_from_rtf(data: bytes) -> list[str]:
    """Extract URLs from RTF attachment text using striprtf."""
    try:
        rtf_text = data.decode("utf-8", errors="replace")
        plain_text = rtf_to_text(rtf_text)
        return _extract_raw_urls(plain_text)
    except Exception:
        logger.warning("Failed to extract URLs from RTF attachment", exc_info=True)
        return []


def _extract_from_plain_bytes(data: bytes) -> list[str]:
    """Decode plain-text attachment bytes and extract URLs."""
    try:
        text = data.decode("utf-8", errors="replace")
        return _extract_raw_urls(text)
    except Exception:
        logger.warning("Failed to extract URLs from plain-text attachment", exc_info=True)
        return []


def _is_pdf_attachment(att: Attachment) -> bool:
    """Return True if the attachment should be parsed as PDF."""
    content_type = (att.content_type or "").lower()
    filename = att.filename.lower()
    return content_type == "application/pdf" or filename.endswith(".pdf")


def _is_docx_attachment(att: Attachment) -> bool:
    """Return True if the attachment should be parsed as DOCX."""
    content_type = (att.content_type or "").lower()
    filename = att.filename.lower()
    return (
        filename.endswith(".docx")
        or content_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )


def _is_xlsx_attachment(att: Attachment) -> bool:
    """Return True if the attachment should be parsed as XLSX."""
    content_type = (att.content_type or "").lower()
    filename = att.filename.lower()
    return (
        filename.endswith(".xlsx")
        or content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )


def _is_pptx_attachment(att: Attachment) -> bool:
    """Return True if the attachment should be parsed as PPTX."""
    content_type = (att.content_type or "").lower()
    filename = att.filename.lower()
    return (
        filename.endswith(".pptx")
        or content_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


def _is_rtf_attachment(att: Attachment) -> bool:
    """Return True if the attachment should be parsed as RTF."""
    content_type = (att.content_type or "").lower()
    filename = att.filename.lower()
    return (
        filename.endswith(".rtf")
        or content_type in {"application/rtf", "text/rtf"}
    )


def _is_plain_text_attachment(att: Attachment) -> bool:
    """Return True if the attachment should be parsed as plain text."""
    content_type = (att.content_type or "").lower()
    filename = att.filename.lower()
    return content_type == "text/plain" or filename.endswith(".txt")


def _is_html_attachment(att: Attachment) -> bool:
    """Return True if the attachment should be parsed as HTML."""
    content_type = (att.content_type or "").lower()
    filename = att.filename.lower()
    return (
        content_type == "text/html"
        or filename.endswith(".html")
        or filename.endswith(".htm")
    )


def _extract_from_attachment(att: Attachment) -> list[str]:
    """Extract URLs from a single attachment based on type."""
    if _is_pdf_attachment(att):
        return _extract_from_pdf(att.data)
    if _is_docx_attachment(att):
        return _extract_from_docx(att.data)
    if _is_xlsx_attachment(att):
        return _extract_from_xlsx(att.data)
    if _is_pptx_attachment(att):
        return _extract_from_pptx(att.data)
    if _is_rtf_attachment(att):
        return _extract_from_rtf(att.data)
    if _is_plain_text_attachment(att):
        return _extract_from_plain_bytes(att.data)
    if _is_html_attachment(att):
        try:
            html = att.data.decode("utf-8", errors="replace")
            return _extract_from_html_content(html)
        except Exception:
            logger.warning(
                "Failed to extract URLs from HTML attachment: %s",
                att.filename,
                exc_info=True,
            )
            return []
    return []


def _make_finding(raw_url: str, source: str, depth: int) -> UrlFinding:
    """Create a UrlFinding with a defanged URL and detection flags."""
    is_shortener = _is_url_shortener(raw_url)
    is_punycode, _ = check_punycode(raw_url)
    is_homoglyph, homoglyph_detail = _check_homoglyphs(raw_url)
    return UrlFinding(
        raw_url=raw_url,
        defanged_url=defang(raw_url),
        source=source,
        depth=depth,
        is_shortener=is_shortener,
        is_punycode=is_punycode,
        homoglyph_detail=homoglyph_detail if is_homoglyph else "",
    )


def _collect_email_urls(
    email: ParsedEmail,
    findings: list[UrlFinding],
    seen: set[tuple[str, str]],
) -> None:
    """Collect URL findings from one email's body and attachments."""

    def add_urls(urls: list[str], source: str, depth: int) -> None:
        for url in urls:
            key = (url.lower(), source)
            if key in seen:
                continue
            seen.add(key)
            findings.append(_make_finding(url, source, depth))

    if email.body_plain:
        add_urls(_extract_raw_urls(email.body_plain), "body:plain", email.depth)

    if email.body_html:
        add_urls(
            _extract_from_html_content(email.body_html),
            "body:html",
            email.depth,
        )

    for attachment in email.attachments:
        source = f"attachment:{attachment.filename}"
        try:
            add_urls(_extract_from_attachment(attachment), source, email.depth)
        except Exception:
            logger.warning(
                "Failed to extract URLs from attachment: %s",
                attachment.filename,
                exc_info=True,
            )


def extract_urls(email: ParsedEmail) -> list[UrlFinding]:
    """Extract, deduplicate, and defang URLs from a parsed email.

    Searches plain and HTML bodies, supported attachment types, and nested
    emails recursively. Deduplication is performed on (raw_url.lower(), source).

    Args:
        email: Parsed email to analyse.

    Returns:
        Deduplicated list of defanged URL findings.
    """
    try:
        findings: list[UrlFinding] = []
        seen: set[tuple[str, str]] = set()

        _collect_email_urls(email, findings, seen)

        for nested in email.nested_emails:
            for finding in extract_urls(nested):
                prefixed_source = f"nested[{nested.depth}]:{finding.source}"
                key = (finding.raw_url.lower(), prefixed_source)
                if key in seen:
                    continue
                seen.add(key)
                findings.append(
                    UrlFinding(
                        raw_url=finding.raw_url,
                        defanged_url=finding.defanged_url,
                        source=prefixed_source,
                        depth=finding.depth,
                        is_shortener=finding.is_shortener,
                        is_punycode=finding.is_punycode,
                        homoglyph_detail=finding.homoglyph_detail,
                    )
                )

        return findings
    except Exception:
        logger.exception("Failed to extract URLs from parsed email")
        return []
